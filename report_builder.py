from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocInches
import pandas as pd
import os
from datetime import datetime

COLORS = {
    'primary': RGBColor(42, 109, 244),
    'primary_hex': '#2A6DF4',
    'secondary': RGBColor(0, 198, 167),
    'secondary_hex': '#00C6A7',
    'warning': RGBColor(245, 166, 35),
    'warning_hex': '#F5A623',
    'text_primary': RGBColor(51, 51, 51),
    'text_secondary': RGBColor(102, 102, 102),
    'bg_light': RGBColor(245, 247, 250),
    'bg_light_hex': '#F5F7FA',
    'border': RGBColor(238, 238, 238),
    'white': RGBColor(255, 255, 255)
}

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def add_gradient_fill(shape, color1, color2):
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0

def add_decorative_elements(slide):
    left = Inches(-0.5)
    top = Inches(0)
    width = Inches(0.8)
    height = Inches(7.5)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.fill.fore_color.brightness = 0.8
    shape.line.fill.background()

def set_font(run, size=18, bold=False, color=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def build_ppt(df, output_dir, output_file="report.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_decorative_elements(slide)

    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "抖音运营月度分析报告"
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=44, bold=True, color=COLORS['primary'])

    date_range = f"{df['日期'].min()} 至 {df['日期'].max()}"
    left = Inches(1)
    top = Inches(3.8)
    width = Inches(8)
    height = Inches(0.6)
    subtitle_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = subtitle_shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = date_range
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=24, color=COLORS['text_secondary'])

    today = datetime.now().strftime("%Y年%m月%d日")
    left = Inches(1)
    top = Inches(6.5)
    width = Inches(8)
    height = Inches(0.5)
    date_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = date_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = f"报告日期：{today}"
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=14, color=COLORS['text_secondary'])

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_decorative_elements(slide)

    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']

    left = Inches(1)
    top = Inches(0.8)
    width = Inches(8)
    height = Inches(0.8)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "目录"
    set_font(p.runs[0], size=32, bold=True, color=COLORS['primary'])

    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(7)
    height = Inches(5)
    content = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content.text_frame
    text_frame.word_wrap = True
    
    items = ["1. 整体概览", "2. 账号详情", "3. 爆款作品", "4. 账号对比", "5. 建议与总结"]
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        set_font(p.runs[0], size=20, color=COLORS['text_primary'])
        p.space_after = Pt(12)

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_decorative_elements(slide)

    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']

    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.6)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "整体概览"
    set_font(p.runs[0], size=32, bold=True, color=COLORS['primary'])

    total_fans = df.sort_values('日期').groupby('账号名称')['粉丝量'].last().sum()
    kpis = [
        ("👥 总粉丝", f"{total_fans:,}"),
        ("📈 总涨粉", f"{df['涨粉量'].sum():,}"),
        ("❤️ 总互动", f"{(df['点赞数'] + df['评论数'] + df['收藏数']).sum():,}")
    ]
    
    for i, (label, value) in enumerate(kpis):
        left = Inches(0.8 + i * 3)
        top = Inches(1.3)
        width = Inches(2.8)
        height = Inches(1.8)
        
        shape = slide.shapes.add_shape(1, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = COLORS['white']
        shape.line.fill.background()
        
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_bottom = Inches(0.15)
        
        p = text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        set_font(p.runs[0], size=16, color=COLORS['text_secondary'])
        
        p = text_frame.add_paragraph()
        p.text = value
        p.alignment = PP_ALIGN.CENTER
        set_font(p.runs[0], size=36, bold=True, color=COLORS['primary'])

    img_path = os.path.join(output_dir, "overview_pie.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1), Inches(3.5), height=Inches(3.5))

    accounts = df['账号名称'].unique()
    for account in accounts:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        add_decorative_elements(slide)

        background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        background.fill.solid()
        background.fill.fore_color.rgb = COLORS['bg_light']

        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(0.6)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"账号详情 - {account}"
        set_font(p.runs[0], size=32, bold=True, color=COLORS['primary'])

        img_path = os.path.join(output_dir, f"detail_{account}.png")
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), height=Inches(5.8))

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_decorative_elements(slide)

    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']

    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.6)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "爆款作品"
    set_font(p.runs[0], size=32, bold=True, color=COLORS['primary'])

    img_path = os.path.join(output_dir, "top_posts.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.3), height=Inches(2.8))

    top_posts = df.sort_values('互动数', ascending=False).head(10)
    table = slide.shapes.add_table(11, 4, Inches(0.5), Inches(4.3), Inches(9), Inches(2.8)).table
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(2)

    headers = ["作品标题", "账号", "互动数", "互动率"]
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = COLORS['white']
                run.font.bold = True
                run.font.size = Pt(14)

    for idx, (_, post) in enumerate(top_posts.iterrows()):
        row = table.rows[idx + 1]
        if idx % 2 == 1:
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['bg_light']
        
        row.cells[0].text = post['作品标题'][:30] + "..."
        row.cells[1].text = post['账号名称']
        row.cells[2].text = f"{post['互动数']:,}"
        row.cells[3].text = f"{post['互动率']:.2%}"
        
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = COLORS['text_secondary']
                    run.font.size = Pt(12)

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_decorative_elements(slide)

    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']

    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.6)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "账号对比"
    set_font(p.runs[0], size=32, bold=True, color=COLORS['primary'])

    img_path = os.path.join(output_dir, "comparison.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.3), height=Inches(5.8))

    left = Inches(1)
    top = Inches(7.2)
    width = Inches(8)
    height = Inches(0.3)
    source_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = source_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "数据来源：抖音后台数据统计"
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=10, color=RGBColor(153, 153, 153))

    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_decorative_elements(slide)

    background = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg_light']

    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(0.6)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "建议与总结"
    set_font(p.runs[0], size=32, bold=True, color=COLORS['primary'])

    left = Inches(0.8)
    top = Inches(1.3)
    width = Inches(8.4)
    height = Inches(5.8)
    content = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content.text_frame
    text_frame.word_wrap = True

    summary_text = """【亮点】
1. 整体粉丝增长趋势良好，各账号均有稳定表现
2. 爆款作品互动率突出，内容质量得到用户认可
3. 账号矩阵布局合理，覆盖多个垂直领域

【问题】
1. 部分账号涨粉波动较大，稳定性有待提升
2. 评论互动率相对较低，需加强用户引导
3. 内容发布频率不均衡，建议优化发布策略

【建议】
1. 针对爆款作品内容特点，持续产出同类型优质内容
2. 增加评论区互动，积极回复用户留言
3. 制定固定发布计划，保持内容更新频率"""

    lines = summary_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = line
        if line.startswith('【'):
            set_font(p.runs[0], size=18, bold=True, color=COLORS['primary'])
            p.space_before = Pt(12)
        else:
            set_font(p.runs[0], size=16, color=COLORS['text_secondary'])
        p.space_after = Pt(6)

    output_path = os.path.join(output_dir, output_file)
    prs.save(output_path)
    return output_path

def build_word(df, output_dir, output_file="report.docx"):
    doc = Document()
    doc.add_heading('抖音运营月度分析报告', 0)
    doc.add_paragraph(f"报告期间：{df['日期'].min()} 至 {df['日期'].max()}")
    
    doc.add_heading('整体概览', level=1)
    total_fans = df.sort_values('日期').groupby('账号名称')['粉丝量'].last().sum()
    doc.add_paragraph(f"总粉丝数：{total_fans:,}")
    doc.add_paragraph(f"总涨粉：{df['涨粉量'].sum():,}")
    doc.add_paragraph(f"总点赞：{df['点赞数'].sum():,}")
    doc.add_paragraph(f"总评论：{df['评论数'].sum():,}")
    doc.add_paragraph(f"总收藏：{df['收藏数'].sum():,}")
    doc.add_paragraph(f"总播放：{df['播放量'].sum():,}")
    
    doc.add_heading('爆款作品', level=1)
    top_posts = df.sort_values('互动数', ascending=False).head(10)
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = '作品标题'
    hdr_cells[1].text = '账号'
    hdr_cells[2].text = '互动数'
    hdr_cells[3].text = '互动率'
    
    for _, post in top_posts.iterrows():
        row_cells = table.add_row().cells
        row_cells[0].text = post['作品标题']
        row_cells[1].text = post['账号名称']
        row_cells[2].text = f"{post['互动数']:,}"
        row_cells[3].text = f"{post['互动率']:.2%}"
    
    doc.add_heading('建议与总结', level=1)
    summary_text = """
【亮点】
1. 整体粉丝增长趋势良好，各账号均有稳定表现
2. 爆款作品互动率突出，内容质量得到用户认可
3. 账号矩阵布局合理，覆盖多个垂直领域

【问题】
1. 部分账号涨粉波动较大，稳定性有待提升
2. 评论互动率相对较低，需加强用户引导
3. 内容发布频率不均衡，建议优化发布策略

【建议】
1. 针对爆款作品内容特点，持续产出同类型优质内容
2. 增加评论区互动，积极回复用户留言
3. 制定固定发布计划，保持内容更新频率
    """
    doc.add_paragraph(summary_text.strip())
    
    output_path = os.path.join(output_dir, output_file)
    doc.save(output_path)
    return output_path
